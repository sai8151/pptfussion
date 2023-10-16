import spacy
import pytextrank
import random
import pdfplumber
import os
from pptx import Presentation
from pptx.util import Inches
# Initialize SpaCy with TextRank
nlp = spacy.load("en_core_web_lg")
nlp.add_pipe("textrank")
def getPoint(text):
    doc = nlp(text)
    summaries = []
    for sent in doc._.textrank.summary(limit_phrases=random.randint(0, 3), limit_sentences=random.randint(2, 3)):
        summaries.append(sent.text)
    return summaries
def create_ppt_slides(pdf_path, output_pptx_path):
    presentation = Presentation()
    with pdfplumber.open(pdf_path) as pdf:
        for page_num in range(len(pdf.pages)):
            page = pdf.pages[page_num]
            page_text = page.extract_text()
            if page_text:
                slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # Blank slide layout
                text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
                text_frame = text_box.text_frame
                summaries = getPoint(page_text)
                for summary in summaries:
                    p = text_frame.add_paragraph()
                    p.text = summary
                presentation.save(output_pptx_path)
    presentation.save(output_pptx_path)
# Set the path to your PDF file and the output PPTX file
pdf_path = "ML2.pdf"
output_pptx_path = "output.pptx"
create_ppt_slides(pdf_path, output_pptx_path)